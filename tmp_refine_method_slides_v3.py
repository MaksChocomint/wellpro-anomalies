from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"d:/diploma/wellpro-anomalies")
SRC = ROOT / "ВКР_Иванов_МС_v2_weights.pptx"
OUT = ROOT / "ВКР_Иванов_МС_v3.pptx"

ASSETS = ROOT / "tmp_ppt_assets"
IMG_SPIKES = ASSETS / "case02_spikes_speed.png"
IMG_STUCK = ASSETS / "case03_stuck_depth_speed.png"
IMG_FFT = ASSETS / "case04_fft_pressure.png"

for p in [SRC, IMG_SPIKES, IMG_STUCK, IMG_FFT]:
    if not p.exists():
        raise FileNotFoundError(p)


# Palette
BLUE_DARK = RGBColor(26, 69, 122)
BLUE = RGBColor(42, 122, 227)
BLUE_LIGHT = RGBColor(235, 244, 255)
BLUE_SOFT = RGBColor(246, 250, 255)
BORDER = RGBColor(198, 216, 242)
TEXT_MAIN = RGBColor(35, 52, 84)
TEXT_MUTED = RGBColor(84, 100, 128)
ORANGE = RGBColor(235, 133, 35)
GREEN = RGBColor(39, 155, 102)
GRAY = RGBColor(148, 156, 170)


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if p.runs:
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def clear_slide(slide):
    for shp in list(slide.shapes):
        shp.element.getparent().remove(shp.element)


def ensure_notes(slide, template_slide):
    ns = slide.notes_slide
    if len(ns.placeholders) == 0:
        for child in list(template_slide.notes_slide.shapes._spTree)[2:]:
            ns.shapes._spTree.insert_element_before(deepcopy(child), "p:extLst")
    return ns


def header(slide, title, subtitle, num):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.78), Inches(0.62), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    add_text(slide, 0.55, 0.32, 11.2, 0.43, title, size=30, bold=True, color=BLUE_DARK)
    add_text(slide, 0.55, 0.84, 12.0, 0.30, subtitle, size=13.2, color=TEXT_MUTED)
    add_text(slide, 12.55, 7.12, 0.45, 0.20, num, size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def footer(slide, text):
    f = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.42), Inches(12.20), Inches(0.62))
    f.fill.solid()
    f.fill.fore_color.rgb = BLUE_LIGHT
    f.line.fill.background()
    add_text(slide, 0.70, 6.53, 11.90, 0.36, text, size=11.8, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, title):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(255, 255, 255)
    c.line.color.rgb = BORDER
    c.line.width = Pt(1.4)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_SOFT
    bar.line.fill.background()
    add_text(slide, x + 0.12, y + 0.10, w - 0.24, 0.24, title, size=15.5, bold=True, color=BLUE_DARK)
    return c


def metric_chip(slide, x, y, label, value, fill=BLUE_LIGHT, value_color=BLUE_DARK):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.20), Inches(0.46))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    add_text(slide, x + 0.07, y + 0.06, 0.44, 0.20, label, size=9.5, bold=True, color=TEXT_MUTED)
    add_text(slide, x + 0.54, y + 0.04, 0.58, 0.26, str(value), size=16, bold=True, color=value_color, align=PP_ALIGN.CENTER)


prs = Presentation(str(SRC))
if len(prs.slides) < 11:
    raise RuntimeError("Expected 11 slides in source deck")

notes_template = prs.slides[0]
s7, s8, s9, s10 = prs.slides[6], prs.slides[7], prs.slides[8], prs.slides[9]

# -------------------- Slide 7 --------------------
clear_slide(s7)
header(
    s7,
    "Три базовых метода: что видно на графике",
    "Без формул: каждый метод ищет свой тип подозрительного поведения.",
    "07",
)

x0, y0, w, h, g = 0.55, 1.30, 3.95, 4.98, 0.23

card(s7, x0, y0, w, h, "Z-score")
s7.shapes.add_picture(str(IMG_SPIKES), Inches(x0 + 0.12), Inches(y0 + 0.55), Inches(w - 0.24), Inches(1.42))
add_text(s7, x0 + 0.14, y0 + 2.04, w - 0.28, 0.40, "Ищет: резкий отрыв от соседних точек.", size=10.8)
add_text(s7, x0 + 0.14, y0 + 2.46, w - 0.28, 0.40, "Вид на графике: одиночный пик/провал.", size=10.8)
metric_chip(s7, x0 + 0.14, y0 + 2.95, "Пример", 34, fill=BLUE_LIGHT, value_color=BLUE)
add_text(s7, x0 + 0.14, y0 + 3.50, w - 0.28, 1.22, "Сильная сторона: хорошо ловит всплески.\nСлабая сторона: может принять смену режима за аномалию.", size=10.3, color=TEXT_MUTED)

x1 = x0 + w + g
card(s7, x1, y0, w, h, "LOF")
s7.shapes.add_picture(str(IMG_SPIKES), Inches(x1 + 0.12), Inches(y0 + 0.55), Inches(w - 0.24), Inches(1.42))
add_text(s7, x1 + 0.14, y0 + 2.04, w - 0.28, 0.40, "Ищет: точку, «чужую» для своего участка.", size=10.8)
add_text(s7, x1 + 0.14, y0 + 2.46, w - 0.28, 0.64, "Просто: рядом кривая ровная, но одна точка выбивается.", size=10.8)
metric_chip(s7, x1 + 0.14, y0 + 2.95, "Пример", 41, fill=BLUE_LIGHT, value_color=BLUE)
add_text(s7, x1 + 0.14, y0 + 3.50, w - 0.28, 1.22, "Сильная сторона: видит локально странные места.\nСлабая сторона: «странно» ≠ «аварийно».", size=10.3, color=TEXT_MUTED)

x2 = x1 + w + g
card(s7, x2, y0, w, h, "FFT")
s7.shapes.add_picture(str(IMG_FFT), Inches(x2 + 0.12), Inches(y0 + 0.55), Inches(w - 0.24), Inches(1.42))
add_text(s7, x2 + 0.14, y0 + 2.04, w - 0.28, 0.40, "Ищет: частые колебания (пульсации).", size=10.8)
add_text(s7, x2 + 0.14, y0 + 2.46, w - 0.28, 0.40, "Вид на графике: «пила» вверх-вниз.", size=10.8)
metric_chip(s7, x2 + 0.14, y0 + 2.95, "Пример", 424, fill=RGBColor(255, 243, 230), value_color=ORANGE)
add_text(s7, x2 + 0.14, y0 + 3.50, w - 0.28, 1.22, "Сильная сторона: лучший для волн и вибраций.\nСлабая сторона: часто перегружает тревогами.", size=10.3, color=TEXT_MUTED)

norm_chip = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.92), Inches(4.35), Inches(0.36))
norm_chip.fill.solid()
norm_chip.fill.fore_color.rgb = RGBColor(243, 248, 255)
norm_chip.line.fill.background()
add_text(s7, 0.68, 5.98, 4.10, 0.22, "Контроль: нормальный файл → Z=0, LOF=0, FFT=0, AMMAD=0", size=9.8, color=BLUE_DARK)

footer(s7, "Срабатывание — это сигнал «проверить участок», а не автоматический диагноз аварии.")

ns7 = ensure_notes(s7, notes_template)
if ns7.notes_text_frame is not None:
    ns7.notes_text_frame.text = (
        "Главная мысль слайда: методы не конкурируют, а смотрят на разные паттерны. "
        "LOF объясняем как «точка не похожа на соседей», без математических терминов."
    )

# -------------------- Slide 8 --------------------
clear_slide(s8)
header(
    s8,
    "Где одиночные методы дают промахи",
    "Два реальных сценария: пропуск события и перегруз тревогами.",
    "08",
)

card(s8, 0.55, 1.30, 6.10, 4.98, "Сценарий A: залипание канала при росте глубины")
s8.shapes.add_picture(str(IMG_STUCK), Inches(0.70), Inches(1.84), Inches(5.80), Inches(2.00))
add_text(s8, 0.72, 3.90, 5.76, 0.54, "Что видим: глубина растет, а скорость «зависла».", size=11.0)
metric_chip(s8, 0.74, 4.48, "Z", 28)
metric_chip(s8, 2.02, 4.48, "LOF", 17)
metric_chip(s8, 3.30, 4.48, "FFT", 0, fill=RGBColor(241, 243, 246), value_color=GRAY)
metric_chip(s8, 4.58, 4.48, "AMMAD", 39, fill=RGBColor(231, 248, 239), value_color=GREEN)
add_text(s8, 0.72, 5.08, 5.76, 1.00, "Вывод: FFT пропускает событие (это не частотная проблема), AMMAD ловит контекст.", size=10.6, color=TEXT_MUTED)

card(s8, 6.78, 1.30, 5.97, 4.98, "Сценарий B: высокочастотные осцилляции")
s8.shapes.add_picture(str(IMG_FFT), Inches(6.94), Inches(1.84), Inches(5.65), Inches(2.00))
add_text(s8, 6.96, 3.90, 5.60, 0.54, "Что видим: резкая «пила» давления.", size=11.0)
metric_chip(s8, 6.98, 4.48, "Z", 27)
metric_chip(s8, 8.26, 4.48, "LOF", 28)
metric_chip(s8, 9.54, 4.48, "FFT", 424, fill=RGBColor(255, 243, 230), value_color=ORANGE)
metric_chip(s8, 10.82, 4.48, "AMMAD", 55, fill=RGBColor(231, 248, 239), value_color=GREEN)
add_text(s8, 6.96, 5.08, 5.60, 1.00, "Вывод: FFT видит слишком много, AMMAD снижает шум для оператора.", size=10.6, color=TEXT_MUTED)

footer(s8, "Одиночный метод полезен точечно, но для мониторинга нужен фильтр с технологическим контекстом.")

ns8 = ensure_notes(s8, notes_template)
if ns8.notes_text_frame is not None:
    ns8.notes_text_frame.text = (
        "Показываем два контрастных кейса: где FFT ничего не видит и где FFT видит слишком много. "
        "Так легче объяснить, зачем нужен гибридный подход."
    )

# -------------------- Slide 9 --------------------
clear_slide(s9)
header(
    s9,
    "Почему AMMAD удобнее инженеру",
    "AMMAD не заменяет решение инженера, а ставит приоритеты в потоке тревог.",
    "09",
)

card(s9, 0.55, 1.30, 6.00, 4.98, "Что делает AMMAD")
add_text(
    s9,
    0.78,
    1.90,
    5.55,
    2.70,
    "1) Объединяет Z-score + LOF + FFT\n"
    "2) Проверяет физические пределы\n"
    "3) Смотрит контекст глубины (бурим/стоим)\n"
    "4) Ловит «залипание» канала\n"
    "5) Требует устойчивость события\n"
    "6) Убирает повторы через cooldown",
    size=11.8,
    color=TEXT_MAIN,
)
add_text(
    s9,
    0.78,
    4.75,
    5.55,
    1.30,
    "Результат: меньше ложной нагрузки и быстрее путь к действительно важным участкам.",
    size=11.2,
    bold=True,
    color=BLUE,
)

card(s9, 6.78, 1.30, 5.97, 2.35, "Быстрый эффект на ваших сценариях")
add_text(s9, 7.02, 1.92, 5.45, 0.45, "Осцилляции: FFT 424 → AMMAD 55", size=13.8, bold=True, color=ORANGE)
add_text(s9, 7.02, 2.38, 5.45, 0.45, "Залипание: FFT 0 → AMMAD 39", size=13.8, bold=True, color=GREEN)
add_text(s9, 7.02, 2.84, 5.45, 0.62, "Смешанный файл: FFT 190 → AMMAD 13", size=13.0, bold=True, color=BLUE_DARK)

card(s9, 6.78, 3.93, 5.97, 2.35, "Что это значит для работы")
add_text(
    s9,
    7.02,
    4.44,
    5.45,
    1.45,
    "• меньше «мусорных» тревог в списке\n"
    "• меньше риска пропустить контекстную проблему\n"
    "• быстрее ручной разбор смены / рейса",
    size=11.4,
    color=TEXT_MAIN,
)

footer(s9, "AMMAD = умный предварительный фильтр. Финальное инженерное решение всегда остается за специалистом.")

ns9 = ensure_notes(s9, notes_template)
if ns9.notes_text_frame is not None:
    ns9.notes_text_frame.text = (
        "Этот слайд даем как управленческий вывод: AMMAD уменьшает шум и не пропускает контекстные сбои, "
        "поэтому экономит время на разбор."
    )

# -------------------- Slide 10 (weights) --------------------
clear_slide(s10)
header(
    s10,
    "Веса AMMAD: почему они разные",
    "Вес показывает, какому типу сигнала больше доверяем для конкретного параметра.",
    "10",
)

card(s10, 0.55, 1.30, 4.05, 1.70, "Как читать тройку весов")
add_text(s10, 0.72, 1.84, 3.70, 0.30, "w_z — скачок, w_lof — локальная странность, w_fft — пульсации", size=10.6, color=TEXT_MAIN)
chip1 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(2.20), Inches(3.70), Inches(0.40))
chip1.fill.solid()
chip1.fill.fore_color.rgb = BLUE_LIGHT
chip1.line.fill.background()
add_text(s10, 0.85, 2.28, 3.45, 0.22, "Пример: 0.2 / 0.4 / 0.4", size=11.2, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

card(s10, 4.78, 1.30, 3.95, 1.70, "Скорость бурения / СПО")
add_text(s10, 4.96, 1.84, 3.58, 0.30, "0.2 / 0.4 / 0.4", size=14.2, bold=True, color=BLUE)
add_text(s10, 4.96, 2.18, 3.58, 0.66, "Много нулей и переходов,\nпоэтому выше роль LOF/FFT.", size=10.4, color=TEXT_MUTED)

card(s10, 8.89, 1.30, 3.86, 1.70, "Давление / расход")
add_text(s10, 9.05, 1.84, 3.52, 0.30, "0.5/0.4/0.1 и 0.6/0.3/0.1", size=12.3, bold=True, color=BLUE)
add_text(s10, 9.05, 2.18, 3.52, 0.66, "FFT ослаблен, чтобы не\nшуметь на штатной пульсации.", size=10.4, color=TEXT_MUTED)

card(s10, 0.55, 3.20, 4.05, 1.92, "Температура / вес на крюке")
add_text(s10, 0.72, 3.74, 3.70, 0.30, "0.8/0.2/0.0 и 0.7/0.2/0.1", size=12.6, bold=True, color=BLUE)
add_text(s10, 0.72, 4.08, 3.70, 0.86, "Медленные и уровневые каналы:\nважнее уровень и выход за диапазон,\nчем частота.", size=10.2, color=TEXT_MUTED)

card(s10, 4.78, 3.20, 3.95, 1.92, "ДМК / уровень в емкости")
add_text(s10, 4.96, 3.74, 3.58, 0.30, "0.3/0.5/0.2 и 0.4/0.5/0.1", size=12.6, bold=True, color=BLUE)
add_text(s10, 4.96, 4.08, 3.58, 0.86, "Чаще важна локальная форма,\nпоэтому у LOF максимальный вес.", size=10.2, color=TEXT_MUTED)

card(s10, 8.89, 3.20, 3.86, 1.92, "Механика: обороты/нагрузка/момент")
add_text(s10, 9.05, 3.70, 3.52, 0.42, "0.4/0.2/0.4\n0.4/0.3/0.3\n0.3/0.4/0.3", size=10.8, bold=True, color=BLUE)
add_text(s10, 9.05, 4.20, 3.52, 0.74, "Смешанная природа: важны\nи скачки, и локальные сбои,\nи колебания.", size=10.0, color=TEXT_MUTED)

footer(s10, "Веса взяты из анализа 15 файлов и помогают уменьшить ложные тревоги на реальных режимах бурения.")

ns10 = ensure_notes(s10, notes_template)
if ns10.notes_text_frame is not None:
    ns10.notes_text_frame.text = (
        "Смысл слайда: веса отражают физику канала. "
        "Они не «магические», а практические — чтобы AMMAD меньше шумел и не пропускал важные паттерны."
    )

prs.save(str(OUT))
print(f"Saved: {OUT}")
